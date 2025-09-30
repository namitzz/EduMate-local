# backend/ingest.py
import os
import hashlib
from pathlib import Path

from chromadb import Client
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer

from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from chunker import split_text
import config

# Supported file types to ingest
SUPPORTED = {".pdf", ".txt", ".md", ".docx", ".pptx", ".html", ".htm"}


def read_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext == ".pptx":
        texts = []
        pres = Presentation(path)
        for slide in pres.slides:
            for shape in slide.shapes:
                # handles text boxes, titles, placeholders, etc.
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    if ext in {".html", ".htm"}:
        html = Path(path).read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        return soup.get_text(" ", strip=True)

    if ext == ".pdf":
        try:
            reader = PdfReader(str(path))
            return "\n".join((page.extract_text() or "") for page in reader.pages)
        except Exception:
            return ""  # fail-soft; skip unreadable pdfs

    # .txt, .md, others we can read as text
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def doc_id_for(path: Path) -> str:
    return hashlib.sha256(str(path).encode()).hexdigest()[:16]


def main():
    # Ensure persistence dir exists
    os.makedirs(config.DATA_DIR, exist_ok=True)

    # Persistent Chroma client (0.5.x persists automatically with persist_directory)
    client = Client(
        Settings(
            persist_directory=str(config.DATA_DIR),
            allow_reset=True,
            is_persistent=True,
        )
    )
    collection = client.get_or_create_collection("edumate")

    # Use the same embedding model here and in retrieval.py
    embedder = SentenceTransformer(config.EMBEDDING_MODEL)

    # Gather files
    files = []
    for p in Path(config.CORPUS_DIR).glob("**/*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED:
            files.append(p)

    print(f"Found {len(files)} files in corpus")
    ids, docs, metas = [], [], []

    for fp in files:
        text = read_text(fp)
        if not text or not text.strip():
            continue

        chunks = split_text(text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)
        for i, ch in enumerate(chunks):
            ids.append(f"{doc_id_for(fp)}-{i}")
            docs.append(ch)
            metas.append({"file": fp.name, "path": str(fp), "chunk": i})

    if not ids:
        print("No content found. Place files in ./corpus and rerun.")
        return

    print(f"Upserting {len(ids)} chunks...")
    embs = embedder.encode(docs, show_progress_bar=True).tolist()
    collection.upsert(ids=ids, documents=docs, embeddings=embs, metadatas=metas)

    # No client.persist() on chromadb 0.5.x — persisted automatically
    print("Ingestion complete. (Chroma at:", config.DATA_DIR, ")")


if __name__ == "__main__":
    main()


